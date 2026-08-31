#!/usr/bin/env python3
"""Build YCompany RIMSS promotion assessment presentation with Nagarro branding."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent
OUT = ROOT / "YCompany-Promotion-Assessment.pptx"
LOGO_DARK = ROOT / "assets" / "nagarro-logo-dark-transparent.png"
LOGO_LIGHT = ROOT / "assets" / "nagarro-logo-light-transparent.png"

# Layout — logo and title spacing
LOGO_LEFT = Inches(0.45)
LOGO_TOP = Inches(0.2)
LOGO_WIDTH = Inches(1.75)
TITLE_LEFT = Inches(3.05)  # clear gap after logo (~0.85")
TITLE_WIDTH = Inches(9.8)

# Nagarro brand (from official logo SVG)
GREEN = RGBColor(0x47, 0xD7, 0xAC)
DARK = RGBColor(0x06, 0x04, 0x1F)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF4, 0xFA, 0xF8)
MID_GRAY = RGBColor(0x66, 0x66, 0x66)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
FOOTER = "Nagarro · YCompany RIMSS · Confidential"
ARROW_COLOR = GREEN
ARROW_GAP = 0.05  # inches between box edge and arrow


def _to_back(slide, shape):
    """Send shape behind later-added content (keeps arrows visible in gaps only)."""
    el = shape._element
    tree = slide.shapes._spTree
    tree.remove(el)
    tree.insert(2, el)


def _arrow_between(slide, x, bottom_y, top_y, send_back=True):
    """Draw a vertical connector in the gap between two diagram boxes."""
    y1 = bottom_y + ARROW_GAP
    y2 = top_y - ARROW_GAP - 0.11
    if y2 <= y1 + 0.04:
        y1 = bottom_y + 0.02
        y2 = top_y - 0.13
    if y2 <= y1:
        return

    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x), Inches(y1), Inches(x), Inches(y2)
    )
    line.line.color.rgb = ARROW_COLOR
    line.line.width = Pt(2.25)

    tri = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ISOSCELES_TRIANGLE,
        Inches(x - 0.09),
        Inches(y2),
        Inches(0.18),
        Inches(0.11),
    )
    tri.fill.solid()
    tri.fill.fore_color.rgb = ARROW_COLOR
    tri.line.fill.background()
    tri.rotation = 180.0

    if send_back:
        _to_back(slide, line)
        _to_back(slide, tri)


def _arrow_down(slide, x, y1, y2):
    """Legacy wrapper — treats y1/y2 as outer box edges."""
    _arrow_between(slide, x, y1, y2)


def _step_label(slide, num, top):
    badge = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL, Inches(0.55), Inches(top), Inches(0.28), Inches(0.28)
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = GREEN
    badge.line.fill.background()
    btf = badge.text_frame
    btf.vertical_anchor = MSO_ANCHOR.MIDDLE
    bp = btf.paragraphs[0]
    bp.alignment = PP_ALIGN.CENTER
    br = bp.add_run()
    br.text = str(num)
    set_run(br, size=10, bold=True, color=DARK)
    return top + 0.28


def _flow_box_height(lines):
    return 0.56 if lines else 0.48


def _draw_flow_steps(slide, steps, *, box_left=0.95, box_width=11.4, arrow_x=6.65):
    """Render numbered flow steps; draw connectors after all boxes (avoids overlap)."""
    bounds: list[tuple[float, float]] = []
    for i, (top, title, lines, fill, title_color, body_color) in enumerate(steps, start=1):
        h = _flow_box_height(lines)
        _step_label(slide, i, top)
        _diagram_box(
            slide, box_left, top, box_width, h, title, lines, fill, title_color, body_color, 10, 8
        )
        bounds.append((top, top + h))

    for i in range(len(bounds) - 1):
        _arrow_between(slide, arrow_x, bounds[i][1], bounds[i + 1][0], send_back=False)


def add_business_context_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_band(slide, "Business Context")

    left_bg = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(1.32), Inches(5.35), Inches(5.35)
    )
    left_bg.fill.solid()
    left_bg.fill.fore_color.rgb = LIGHT_BG
    left_bg.line.color.rgb = MID_GRAY

    right_bg = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(7.45), Inches(1.32), Inches(5.35), Inches(5.35)
    )
    right_bg.fill.solid()
    right_bg.fill.fore_color.rgb = LIGHT_BG
    right_bg.line.color.rgb = GREEN

    _section_header(slide, 0.7, 1.42, 5.05, "Before — Legacy Storefront", fill=DARK, text_color=WHITE)
    _section_header(slide, 7.6, 1.42, 5.05, "After — RIMSS (As Built)", fill=GREEN, text_color=DARK)

    challenges = [
        ("Performance", "First paint > 60 s on retail Wi‑Fi · UI lag during checkout"),
        ("Search", "Slow text search · no filters · no relevance ranking"),
        ("Inventory", "Out-of-stock SKUs could still be added to cart"),
        ("Experience", "Inconsistent layout across mobile and desktop"),
    ]
    outcomes = [
        ("Performance", "FCP target < 3 s · RSC · next/image · route loading states"),
        ("Search", "Elastic Cloud + fallback · brand/price/sort filters"),
        ("Inventory", "clampAddQuantity() enforces per-SKU stock in cart"),
        ("Experience", "Responsive MUI 9 · drawer nav · mobile cart cards"),
    ]

    for i, ((ct, cd), (ot, od)) in enumerate(zip(challenges, outcomes)):
        top = 2.0 + i * 1.05
        _bullet_card(slide, 0.7, top, 5.05, 0.92, ct, [cd], header_fill=DARK)
        _bullet_card(slide, 7.6, top, 5.05, 0.92, ot, [od])

    # Centre transform arrow
    mid = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.CHEVRON, Inches(6.15), Inches(3.35), Inches(0.95), Inches(1.2)
    )
    mid.fill.solid()
    mid.fill.fore_color.rgb = GREEN
    mid.line.fill.background()

    transform = slide.shapes.add_textbox(Inches(5.95), Inches(2.55), Inches(1.35), Inches(0.55))
    tfp = transform.text_frame.paragraphs[0]
    tfp.alignment = PP_ALIGN.CENTER
    tr = tfp.add_run()
    tr.text = "Modern\nre-platform"
    set_run(tr, size=11, bold=True, color=DARK)

    add_bullets(
        slide,
        [
            "RIMSS replaces legacy pain points with a demo-ready Next.js storefront and a clear path to PostgreSQL production data.",
        ],
        top=6.55,
        height=0.35,
        size=11,
    )
    add_footer(slide)



def set_run(run, *, size=18, bold=False, color=DARK, name="Calibri"):
    run.font.name = name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def _apply_list_style(paragraph, *, numbered: bool = False):
    p_pr = paragraph._p.get_or_add_pPr()
    for tag in ("a:buNone", "a:buChar", "a:buAutoNum", "a:buBlip"):
        existing = p_pr.find(qn(tag))
        if existing is not None:
            p_pr.remove(existing)
    if numbered:
        bullet = p_pr.makeelement(qn("a:buAutoNum"), {"type": "arabicPeriod"})
    else:
        bullet = p_pr.makeelement(qn("a:buChar"), {"char": "•"})
    p_pr.insert(0, bullet)
    paragraph.level = 0


def add_footer(slide, text: str = FOOTER):
    bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), SLIDE_H - Inches(0.35), SLIDE_W, Inches(0.35)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = DARK
    bar.line.fill.background()
    tf = bar.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = text
    set_run(run, size=10, color=WHITE)


def add_logo(slide, *, variant: str = "dark", left=LOGO_LEFT, top=LOGO_TOP):
    logo = LOGO_LIGHT if variant == "light" else LOGO_DARK
    if logo.exists():
        slide.shapes.add_picture(str(logo), left, top, width=LOGO_WIDTH)


def add_title_band(slide, title: str, subtitle: str | None = None):
    band = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, Inches(1.15)
    )
    band.fill.solid()
    band.fill.fore_color.rgb = WHITE
    band.line.fill.background()

    accent = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(1.12), SLIDE_W, Inches(0.06)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = GREEN
    accent.line.fill.background()

    add_logo(slide, variant="dark")

    box = slide.shapes.add_textbox(TITLE_LEFT, Inches(0.28), TITLE_WIDTH, Inches(0.75))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    set_run(run, size=28, bold=True, color=DARK)

    if subtitle:
        sub = slide.shapes.add_textbox(TITLE_LEFT, Inches(0.78), TITLE_WIDTH, Inches(0.35))
        sp = sub.text_frame.paragraphs[0]
        sr = sp.add_run()
        sr.text = subtitle
        set_run(sr, size=14, color=MID_GRAY)


def add_bullets(
    slide,
    items: list[str],
    left=0.7,
    top=1.45,
    width=11.8,
    height=5.5,
    size=17,
    numbered: bool = False,
    color: RGBColor = DARK,
    space_after: Pt = Pt(8),
):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for i, item in enumerate(items):
        if not item.strip():
            continue
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = space_after
        _apply_list_style(p, numbered=numbered)
        run = p.add_run()
        run.text = item
        set_run(run, size=size, color=color)
    return box


def add_two_column(
    slide,
    slide_title: str,
    left_title: str,
    left_items: list[str],
    right_title: str,
    right_items: list[str],
):
    add_title_band(slide, slide_title)

    lt = slide.shapes.add_textbox(Inches(0.7), Inches(1.35), Inches(5.8), Inches(0.4))
    ltp = lt.text_frame.paragraphs[0]
    ltr = ltp.add_run()
    ltr.text = left_title
    set_run(ltr, size=20, bold=True, color=GREEN)

    rt = slide.shapes.add_textbox(Inches(6.85), Inches(1.35), Inches(5.8), Inches(0.4))
    rtp = rt.text_frame.paragraphs[0]
    rtr = rtp.add_run()
    rtr.text = right_title
    set_run(rtr, size=20, bold=True, color=GREEN)

    add_bullets(slide, left_items, left=0.7, top=1.85, width=5.8, height=5.0, size=16)
    add_bullets(slide, right_items, left=6.85, top=1.85, width=5.8, height=5.0, size=16)
    add_footer(slide)


def add_table_slide(slide, title, headers, rows):
    add_title_band(slide, title)
    cols, row_count = len(headers), len(rows) + 1
    table_shape = slide.shapes.add_table(
        row_count, cols, Inches(0.7), Inches(1.55), Inches(11.9), Inches(0.42 * row_count)
    )
    table = table_shape.table
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = GREEN
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                set_run(r, size=13, bold=True, color=DARK)
    for r_idx, row in enumerate(rows, start=1):
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = val
            if r_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_BG
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    set_run(run, size=12, color=DARK)
    add_footer(slide)


def _diagram_box(slide, left, top, width, height, title, lines, fill, title_color, body_color, title_size=11, body_size=9):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = GREEN
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.06)
    p0 = tf.paragraphs[0]
    p0.alignment = PP_ALIGN.CENTER
    r0 = p0.add_run()
    r0.text = title
    set_run(r0, size=title_size, bold=True, color=title_color)
    for line in lines:
        p = tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = line
        set_run(r, size=body_size, color=body_color)
    return shape


def _section_header(slide, left, top, width, title, fill=GREEN, text_color=DARK):
    bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(0.42)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = fill
    bar.line.fill.background()
    tf = bar.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.12)
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    set_run(r, size=14, bold=True, color=text_color)
    return bar


def _bullet_card(slide, left, top, width, height, title, items, header_fill=GREEN, body_fill=WHITE):
    card = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = body_fill
    card.line.color.rgb = GREEN
    card.line.width = Pt(1.25)

    hdr = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(0.38),
    )
    hdr.fill.solid()
    hdr.fill.fore_color.rgb = header_fill
    hdr.line.fill.background()

    htf = hdr.text_frame
    htf.vertical_anchor = MSO_ANCHOR.MIDDLE
    htf.margin_left = Inches(0.1)
    hp = htf.paragraphs[0]
    hr = hp.add_run()
    hr.text = title
    set_run(hr, size=11, bold=True, color=DARK if header_fill == GREEN else WHITE)

    body = slide.shapes.add_textbox(
        Inches(left + 0.12), Inches(top + 0.42), Inches(width - 0.24), Inches(height - 0.5)
    )
    btf = body.text_frame
    btf.word_wrap = True
    for i, item in enumerate(items):
        p = btf.paragraphs[0] if i == 0 else btf.add_paragraph()
        p.space_after = Pt(4)
        _apply_list_style(p, numbered=False)
        br = p.add_run()
        br.text = item
        set_run(br, size=9, color=DARK)
    return card


def add_notifications_quality_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_band(slide, "Notifications & Quality")

    # Column backgrounds
    left_bg = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(1.32), Inches(6.05), Inches(5.15)
    )
    left_bg.fill.solid()
    left_bg.fill.fore_color.rgb = LIGHT_BG
    left_bg.line.color.rgb = GREEN

    right_bg = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(6.75), Inches(1.32), Inches(6.05), Inches(5.15)
    )
    right_bg.fill.solid()
    right_bg.fill.fore_color.rgb = LIGHT_BG
    right_bg.line.color.rgb = GREEN

    _section_header(slide, 0.7, 1.42, 5.75, "Transactional Email — Resend")

    _diagram_box(
        slide,
        0.75,
        1.95,
        5.65,
        0.72,
        "Email layer — src/lib/email/",
        [
            "send.ts → Resend REST API when RESEND_API_KEY is set",
            "templates.ts → HTML + plain-text · console fallback for local dev",
        ],
        WHITE,
        DARK,
        DARK,
        11,
        9,
    )

    _bullet_card(
        slide,
        0.75,
        2.82,
        1.75,
        1.55,
        "Welcome",
        ["Trigger: POST /api/auth/signup", "New user onboarding", "sendWelcomeEmail()"],
    )
    _bullet_card(
        slide,
        2.65,
        2.82,
        1.75,
        1.55,
        "Order confirm",
        ["Trigger: Stripe webhook", "+ POST /api/orders/sync", "Idempotent delivery"],
    )
    _bullet_card(
        slide,
        4.55,
        2.82,
        1.75,
        1.55,
        "Cancellation",
        ["Trigger: POST …/cancel", "status → cancelled", "sendOrderCancellationEmail()"],
    )

    _arrow_between(slide, 3.62, 4.37, 4.58, send_back=False)
    _diagram_box(
        slide,
        0.75,
        4.58,
        5.65,
        0.55,
        "Resend → verified domain · onboarding@resend.dev in dev",
        [],
        GREEN,
        DARK,
        DARK,
        10,
        9,
    )

    _section_header(slide, 6.9, 1.42, 5.75, "Observability — env-gated no-ops", fill=DARK, text_color=WHITE)

    _bullet_card(
        slide,
        6.95,
        1.95,
        1.75,
        1.35,
        "Sentry",
        ["Errors + stack traces", "Session replay", "Performance traces"],
        header_fill=DARK,
    )
    _bullet_card(
        slide,
        8.85,
        1.95,
        1.75,
        1.35,
        "Google GA4",
        ["add_to_cart · purchase", "begin_checkout · search", "apply_coupon events"],
        header_fill=DARK,
    )
    _bullet_card(
        slide,
        10.75,
        1.95,
        1.75,
        1.35,
        "New Relic",
        ["Server APM agent", "Docker / K8s preload", "Distributed tracing"],
        header_fill=DARK,
    )

    _section_header(slide, 6.9, 3.45, 5.75, "Quality Gates — CI enforced")

    metrics = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(6.95), Inches(3.98), Inches(5.55), Inches(1.15)
    )
    metrics.fill.solid()
    metrics.fill.fore_color.rgb = WHITE
    metrics.line.color.rgb = GREEN

    mtf = metrics.text_frame
    mtf.margin_left = Inches(0.15)
    mtf.margin_top = Inches(0.1)
    metric_rows = [
        ("Vitest 4 + React Testing Library", "201 automated tests in reference impl"),
        ("Coverage thresholds", "Lines / functions ≥ 90% · branches ≥ 85%"),
        ("Static analysis", "ESLint (eslint-config-next) · TypeScript strict"),
    ]
    for i, (label, detail) in enumerate(metric_rows):
        p = mtf.paragraphs[0] if i == 0 else mtf.add_paragraph()
        p.space_after = Pt(6)
        _apply_list_style(p, numbered=True)
        r = p.add_run()
        r.text = f"{label} — {detail}"
        set_run(r, size=10, bold=(i == 0), color=DARK)

    _diagram_box(
        slide,
        6.95,
        5.28,
        5.55,
        0.62,
        "Pattern: all integrations optional until env vars are configured",
        ["Frictionless local dev · staged production rollout · reportError() central hook"],
        GREEN,
        DARK,
        DARK,
        10,
        8,
    )

    add_bullets(
        slide,
        [
            "Emails and observability share the same graceful-degradation model — demos work without cloud credentials.",
            "Order emails tie into checkout and cancel flows (see Order Flow slide); GA4 mirrors the full commerce funnel.",
        ],
        top=6.55,
        height=0.4,
        size=10,
    )
    add_footer(slide)


def add_search_flow_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_band(slide, "Search Architecture — Flow")

    _diagram_box(
        slide, 6.55, 1.30, 5.8, 0.50,
        "Header autocomplete (parallel)",
        ["ProductSearchAutocomplete · in-memory only"],
        LIGHT_BG, DARK, DARK, 10, 8,
    )

    main_steps = [
        (1.30, "User on /search — query + brand / price / sort filters", [], DARK, WHITE, WHITE),
        (1.98, "SearchClient — 300 ms debounce · useProductSearch (TanStack Query 5)", ["URL sync via router.replace"], WHITE, DARK, DARK),
        (2.72, "GET /api/products/search?q=&brands=&min=&max=&sort=", [], GREEN, DARK, DARK),
        (3.40, "API guards — rateLimit 120/min/IP · parseFiltersFromSearchParams (Zod)", [], WHITE, DARK, DARK),
        (4.08, "searchProducts() facade — src/lib/search.ts", [], DARK, WHITE, WHITE),
    ]
    bounds: list[tuple[float, float]] = []
    for i, (top, title, lines, fill, tc, bc) in enumerate(main_steps, start=1):
        h = _flow_box_height(lines)
        _step_label(slide, i, top)
        _diagram_box(slide, 0.95, top, 11.4, h, title, lines, fill, tc, bc, 10, 8)
        bounds.append((top, top + h))

    split_top = 4.82
    split_h = 0.90
    _step_label(slide, 6, split_top)
    _diagram_box(slide, 0.95, split_top, 5.15, split_h, "Elastic Cloud (primary)",
        ["Index ycompany-products", "multi_match + filters", "npm run search:index"], GREEN, DARK, DARK, 10, 8)
    _diagram_box(slide, 7.20, split_top, 5.15, split_h, "In-memory fallback",
        ["When ES unset / fails", "searchCatalog()", "Zero-config demos"], LIGHT_BG, DARK, DARK, 10, 8)

    resp_top = 5.88
    resp_h = 0.54
    _step_label(slide, 7, resp_top)
    _diagram_box(slide, 1.40, resp_top, 10.55, resp_h,
        "Response { products, total, tookMs, source } · Cache-Control 30s / SWR 60",
        [], WHITE, DARK, DARK, 10, 8)

    ax = 6.65
    for i in range(len(bounds) - 1):
        _arrow_between(slide, ax, bounds[i][1], bounds[i + 1][0], send_back=False)
    _arrow_between(slide, ax, bounds[-1][1], split_top, send_back=False)
    _arrow_between(slide, 3.02, split_top + split_h, resp_top, send_back=False)
    _arrow_between(slide, 9.27, split_top + split_h, resp_top, send_back=False)

    add_bullets(slide, [
        "Elastic Cloud optional — set ELASTICSEARCH_CLOUD_ID + ELASTICSEARCH_API_KEY, then npm run search:index.",
        "Header autocomplete bypasses the API for instant typeahead; full search uses the flow above.",
    ], top=6.55, height=0.40, size=10)
    add_footer(slide)


def add_checkout_flow_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_band(slide, "Checkout Flow")

    steps = [
        (1.28, "Browse catalog — Server Components render PDP / category (RSC + next/image)", [], DARK, WHITE, WHITE),
        (1.96, "Add to cart — CartContext + localStorage · clampAddQuantity()", ["GA4: add_to_cart"], WHITE, DARK, DARK),
        (2.64, "Apply coupon (optional) — CouponField → POST /api/coupons/validate", ["WELCOME10 · SAVE20 · FREESHIP"], LIGHT_BG, DARK, DARK),
        (3.32, "Click Checkout — useCheckout → POST /api/checkout (rate-limited)", ["GA4: begin_checkout"], WHITE, DARK, DARK),
        (4.00, "Stripe Checkout Session — shipping · coupon mirrored", ["stripe.checkout.sessions.create()"], GREEN, DARK, DARK),
        (4.68, "Redirect to Stripe-hosted UI", ["Apple Pay · Google Pay · 3DS · PCI on Stripe"], WHITE, DARK, DARK),
        (5.36, "Customer completes payment", ["checkout.session.completed queued"], DARK, WHITE, WHITE),
    ]
    _draw_flow_steps(slide, steps)

    add_bullets(slide, [
        "Hosted Stripe Checkout keeps YCompany out of PCI scope.",
        "Zod-validated payloads; inventory clamping prevents overselling before payment.",
    ], top=6.05, height=0.40, size=10)
    add_footer(slide)


def add_order_flow_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_band(slide, "Order Flow & Lifecycle")

    ax = 6.65
    b1_top, b1_h = 1.28, 0.48
    _step_label(slide, 1, b1_top)
    _diagram_box(slide, 0.95, b1_top, 11.4, b1_h,
        "Stripe checkout.session.completed (signed webhook)", [], DARK, WHITE, WHITE, 10, 8)

    b2_top, b2_h = 1.96, 0.58
    _step_label(slide, 2, b2_top)
    _diagram_box(slide, 0.95, b2_top, 11.4, b2_h,
        "POST /api/webhooks/stripe — verify signature · idempotent handler",
        ["createOrderFromCheckout() keyed on stripeSessionId"], GREEN, DARK, DARK, 10, 8)

    b3_top, b3_h = 2.70, 0.72
    _step_label(slide, 3, b3_top)
    _diagram_box(slide, 0.95, b3_top, 5.35, b3_h, "Order persisted", ["JSON store · status: processing"], WHITE, DARK, DARK, 10, 8)
    _diagram_box(slide, 7.0, b3_top, 5.35, b3_h, "Notify & track", ["Resend confirmation · GA4 purchase"], WHITE, DARK, DARK, 10, 8)

    b4_top, b4_h = 3.58, 0.58
    _step_label(slide, 4, b4_top)
    _diagram_box(slide, 0.95, b4_top, 11.4, b4_h,
        "/checkout/success — useSyncOrder fallback → POST /api/orders/sync",
        ["Reconciles when webhooks unavailable in local dev"], LIGHT_BG, DARK, DARK, 10, 8)

    b5_top, b5_h = 4.30, 0.56
    _step_label(slide, 5, b5_top)
    _diagram_box(slide, 0.95, b5_top, 11.4, b5_h,
        "Account — GET /api/orders · GET /api/orders/[id] (middleware auth)",
        ["/account dashboard · order history"], WHITE, DARK, DARK, 10, 8)

    b6_top, b6_h = 4.98, 0.58
    _step_label(slide, 6, b6_top)
    _diagram_box(slide, 0.95, b6_top, 11.4, b6_h,
        "Cancel while processing — POST /api/orders/[id]/cancel",
        ["status → cancelled · Resend cancellation email"], WHITE, DARK, DARK, 10, 8)

    b7_top, b7_h = 5.72, 0.46
    _step_label(slide, 7, b7_top)
    life = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.95), Inches(b7_top), Inches(11.4), Inches(b7_h)
    )
    life.fill.solid()
    life.fill.fore_color.rgb = GREEN
    life.line.color.rgb = DARK
    ltf = life.text_frame
    ltf.vertical_anchor = MSO_ANCHOR.MIDDLE
    lp = ltf.paragraphs[0]
    lp.alignment = PP_ALIGN.CENTER
    lr = lp.add_run()
    lr.text = "Lifecycle: processing → cancelled | shipped | delivered"
    set_run(lr, size=11, bold=True, color=DARK)

    _arrow_between(slide, ax, b1_top + b1_h, b2_top, send_back=False)
    _arrow_between(slide, ax, b2_top + b2_h, b3_top, send_back=False)
    _arrow_between(slide, ax, b3_top + b3_h, b4_top, send_back=False)
    _arrow_between(slide, ax, b4_top + b4_h, b5_top, send_back=False)
    _arrow_between(slide, ax, b5_top + b5_h, b6_top, send_back=False)
    _arrow_between(slide, ax, b6_top + b6_h, b7_top, send_back=False)

    add_bullets(slide, [
        "Webhook and sync paths are idempotent — safe to retry without duplicate orders or emails.",
        "Cancellation is user-scoped and only allowed while status is processing.",
    ], top=6.30, height=0.38, size=10)
    add_footer(slide)


def add_frontend_architecture_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_band(slide, "Frontend Architecture (High Level)")

    b0_top, b0_h = 1.35, 0.52
    _diagram_box(slide, 3.8, b0_top, 5.7, b0_h, "Browser — React 19 + Next.js 16 App Router", [], DARK, WHITE, WHITE, 12, 9)

    row_top, row_h = 2.15, 1.02
    _diagram_box(slide, 0.55, row_top, 3.85, row_h, "Server Components (RSC)",
        ["Home · Products · PDP", "Server HTML · SEO"], LIGHT_BG, DARK, DARK)
    _diagram_box(slide, 4.75, row_top, 3.85, row_h, "Client Components",
        ["SearchClient · Cart · Auth", "CouponField · Account"], WHITE, DARK, DARK)
    _diagram_box(slide, 8.95, row_top, 3.85, row_h, "Shared Shell",
        ["Header / Footer", "middleware · loading.tsx"], LIGHT_BG, DARK, DARK)

    p_top, p_h = 3.38, 0.68
    _diagram_box(slide, 1.0, p_top, 11.35, p_h,
        "Presentation — MUI 9 · Emotion SSR · AppRouterCacheProvider · ThemeProvider",
        [], GREEN, DARK, DARK, 11, 9)

    s_top, s_h = 4.28, 0.92
    _diagram_box(slide, 1.0, s_top, 5.5, s_h, "Client State & Data",
        ["TanStack Query 5", "CartContext + localStorage", "useProductSearch"], WHITE, DARK, DARK)
    _diagram_box(slide, 6.85, s_top, 5.5, s_h, "Cross-cutting UI",
        ["ProductSearchAutocomplete", "GA / Sentry / New Relic", "i18n helpers"], WHITE, DARK, DARK)

    a_top, a_h = 5.42, 0.58
    _diagram_box(slide, 2.2, a_top, 8.95, a_h,
        "API Client — fetch /api/* · Zod-validated DTOs · shared query keys",
        [], DARK, WHITE, WHITE, 11, 9)

    _arrow_between(slide, 6.65, b0_top + b0_h, row_top, send_back=False)
    for x in (2.45, 6.65, 10.85):
        _arrow_between(slide, x, row_top + row_h, p_top, send_back=False)
    _arrow_between(slide, 6.65, p_top + p_h, s_top, send_back=False)
    _arrow_between(slide, 6.65, s_top + s_h, a_top, send_back=False)

    add_bullets(slide, [
        "RSC delivers fast first paint; client islands handle interactivity only where needed.",
        "TanStack Query caches server data; CartContext owns cart persistence and stock clamping.",
    ], top=6.15, height=0.45, size=11)
    add_footer(slide)


def add_system_architecture_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_band(slide, "System Architecture (High Level)")

    b1_top, b1_h = 1.32, 0.46
    _diagram_box(slide, 4.0, b1_top, 5.3, b1_h, "Customer Browser", [], DARK, WHITE, WHITE, 12, 9)

    b2_top, b2_h = 1.98, 0.46
    _diagram_box(slide, 3.2, b2_top, 6.95, b2_h,
        "Vercel Edge / CDN — security headers · image optimization",
        [], GREEN, DARK, DARK, 11, 9)

    c_top, c_h = 2.62, 2.05
    container = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(c_top), Inches(12.25), Inches(c_h)
    )
    container.fill.solid()
    container.fill.fore_color.rgb = LIGHT_BG
    container.line.color.rgb = DARK
    label = slide.shapes.add_textbox(Inches(0.75), Inches(c_top + 0.06), Inches(11.9), Inches(0.28))
    lr = label.text_frame.paragraphs[0].add_run()
    lr.text = "Next.js 16 Monolith — React 19 · TypeScript 5 · App Router"
    set_run(lr, size=12, bold=True, color=DARK)

    inner_top = c_top + 0.38
    inner_h = 1.32
    _diagram_box(slide, 0.75, inner_top, 3.55, inner_h, "Server Components",
        ["Home · PDP · Search", "MUI + TanStack Q"], WHITE, DARK, DARK, 10, 8)
    _diagram_box(slide, 4.55, inner_top, 3.85, inner_h, "Route Handlers (/api/*)",
        ["checkout · orders · search", "auth · webhooks"], WHITE, DARK, DARK, 10, 8)
    _diagram_box(slide, 9.65, inner_top, 2.85, inner_h, "middleware.ts",
        ["Protect /account", "NextAuth session"], WHITE, DARK, DARK, 10, 8)

    svc_top, svc_h = 4.88, 0.50
    _diagram_box(slide, 2.0, svc_top, 9.35, svc_h,
        "Domain Services — orders · coupons · search · inventory · stripe · email (Resend)",
        [], GREEN, DARK, DARK, 11, 9)

    int_top, int_h = 5.58, 0.70
    _diagram_box(slide, 0.7, int_top, 2.85, int_h, "NextAuth v5", ["Credentials + OAuth"], WHITE, DARK, DARK, 10, 8)
    _diagram_box(slide, 3.75, int_top, 2.85, int_h, "Stripe API", ["Checkout · webhooks"], WHITE, DARK, DARK, 10, 8)
    _diagram_box(slide, 6.8, int_top, 2.85, int_h, "Elastic Cloud", ["Product index · fallback"], WHITE, DARK, DARK, 10, 8)
    _diagram_box(slide, 9.85, int_top, 2.85, int_h, "Persistence", ["JSON demo · PostgreSQL next"], WHITE, DARK, DARK, 10, 8)

    ax = 6.65
    _arrow_between(slide, ax, b1_top + b1_h, b2_top, send_back=False)
    _arrow_between(slide, ax, b2_top + b2_h, c_top, send_back=False)
    _arrow_between(slide, ax, c_top + c_h, svc_top, send_back=False)
    _arrow_between(slide, ax, svc_top + svc_h, int_top, send_back=False)

    note = slide.shapes.add_textbox(Inches(0.7), Inches(6.38), Inches(11.9), Inches(0.40))
    np = note.text_frame.paragraphs[0]
    _apply_list_style(np, numbered=False)
    nr = np.add_run()
    nr.text = "Cross-cutting: Sentry · GA4 · New Relic APM · Resend · scalable on Vercel or Kubernetes"
    set_run(nr, size=10, color=MID_GRAY)
    add_footer(slide)


def add_summary_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_band(slide, "Summary — Engineering Value")

    metrics = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(1.32), Inches(12.25), Inches(0.72)
    )
    metrics.fill.solid()
    metrics.fill.fore_color.rgb = DARK
    metrics.line.fill.background()
    mtf = metrics.text_frame
    mtf.vertical_anchor = MSO_ANCHOR.MIDDLE
    mp = mtf.paragraphs[0]
    mp.alignment = PP_ALIGN.CENTER
    mr = mp.add_run()
    mr.text = "201 tests  ·  ~95% line coverage  ·  14 storefront screens  ·  Next.js 16 monolith  ·  demo-ready today"
    set_run(mr, size=12, bold=True, color=GREEN)

    pillars = [
        ("End-to-end ownership", "Architecture · implementation · Vitest suite · DAR / Solution Approach · k8s + Docker manifests"),
        ("Modern UX stack", "Next.js 16 · React 19 · MUI 9 · accessible responsive layouts · PCI-safe Stripe Checkout"),
        ("Resilient integrations", "Elastic search + in-memory fallback · Resend email + console fallback · optional observability"),
        ("Production patterns", "Zod validation · rate limiting · signed Stripe webhooks · idempotent orders & emails"),
        ("Quality gates", "Vitest + RTL · coverage ≥ 90% lines/functions · ESLint · TypeScript strict in CI"),
        ("Future-ready", "Clear PostgreSQL migration path · same storefront layer · enterprise K8s portability"),
    ]

    for i, (title, detail) in enumerate(pillars):
        col = i % 3
        row = i // 3
        left = 0.70 + col * 4.05
        top = 2.22 + row * 1.95
        _bullet_card(slide, left, top, 3.85, 1.75, title, [detail])

    outcome = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(6.18), Inches(12.25), Inches(0.52)
    )
    outcome.fill.solid()
    outcome.fill.fore_color.rgb = GREEN
    outcome.line.color.rgb = DARK
    otf = outcome.text_frame
    otf.vertical_anchor = MSO_ANCHOR.MIDDLE
    op = otf.paragraphs[0]
    op.alignment = PP_ALIGN.CENTER
    orun = op.add_run()
    orun.text = (
        "Delivered a promotion-ready reference implementation — fast demos locally, staged cloud rollout, enterprise deploy option"
    )
    set_run(orun, size=11, bold=True, color=DARK)

    add_footer(slide)


def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # 1 — Title
    s1 = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s1.background.fill
    bg.solid()
    bg.fore_color.rgb = DARK
    add_logo(s1, variant="light", top=Inches(0.55), left=Inches(0.7))

    accent = s1.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(2.55), SLIDE_W, Inches(0.08)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = GREEN
    accent.line.fill.background()

    title = s1.shapes.add_textbox(Inches(0.7), Inches(2.85), Inches(11.9), Inches(1.2))
    tp = title.text_frame.paragraphs[0]
    tr = tp.add_run()
    tr.text = "Promotion Assessment"
    set_run(tr, size=40, bold=True, color=WHITE)

    sub = s1.shapes.add_textbox(Inches(0.7), Inches(3.85), Inches(11.9), Inches(1.5))
    stf = sub.text_frame
    lines = [
        "YCompany RIMSS — Retail Inventory Management Software System",
        "Tech Stack & Architecture Overview",
        "Presenter: Pawan Gupta · Nagarro Software Pvt. Ltd.",
        "Date: August 2026",
    ]
    for i, line in enumerate(lines):
        p = stf.paragraphs[0] if i == 0 else stf.add_paragraph()
        p.space_after = Pt(6)
        _apply_list_style(p, numbered=True)
        r = p.add_run()
        r.text = line
        set_run(r, size=16 if i < 2 else 14, color=GREEN if i == 0 else WHITE)

    # 2 — Agenda
    s2 = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_band(s2, "Agenda")
    add_bullets(
        s2,
        [
            "Project context and business objectives",
            "Technology stack (as implemented)",
            "Architecture flow diagrams — frontend, system, checkout, order, search",
            "Backend APIs, notifications, and observability",
            "Summary of engineering outcomes",
        ],
        numbered=True,
        space_after=Pt(18),
    )
    add_footer(s2)

    # 3 — Project overview
    s3 = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_band(s3, "Project Overview")
    add_bullets(
        s3,
        [
            "Client: YCompany — luxury countryside fashion brand (RIMSS programme).",
            "Deliverable: Modern B2C ecommerce storefront with inventory-aware cart and order management.",
            "Reference implementation: Next.js 16 monolith — App Router pages + Route Handlers in one codebase.",
            "Scope: Catalog browse & search, PDP, cart, Stripe checkout, account & orders, coupons, auth.",
            "Demo-ready today; PostgreSQL migration planned as the production data-layer next step.",
            "Documents: Solution Approach, DAR, and Estimates aligned to the as-built stack.",
        ],
        space_after=Pt(18),
    )
    add_footer(s3)

    # 4 — Business context (visual)
    add_business_context_slide(prs)

    # 5 — Tech stack table
    s5 = prs.slides.add_slide(prs.slide_layouts[6])
    add_table_slide(
        s5,
        "Technology Stack (As Implemented)",
        ["Layer", "Technology", "Role"],
        [
            ["Frontend", "Next.js 16 · React 19 · TypeScript 5", "App Router, RSC, SSR/CSR hybrid"],
            ["UI", "Material UI 9 · Emotion 11", "Accessible components, SSR cache provider"],
            ["Client data", "TanStack Query 5", "Cached API calls, checkout & order mutations"],
            ["Auth", "NextAuth v5", "Credentials + OAuth (Google, GitHub, Facebook, Apple)"],
            ["Payments", "Stripe Checkout", "Hosted PCI scope, webhooks, coupons"],
            ["Search", "Elastic Cloud + in-memory fallback", "searchProducts() facade, npm run search:index"],
            ["Email", "Resend REST API", "Welcome, order confirm, cancellation emails"],
            ["Observability", "Sentry · GA4 · New Relic", "Errors, analytics, APM (env-gated no-ops)"],
            ["Testing", "Vitest 4 · RTL · jsdom", "201 tests · ~95% line coverage in CI"],
            ["Deploy", "Vercel · Docker · Kubernetes", "Primary cloud + portable enterprise manifests"],
        ],
    )

    # 6 — Frontend architecture diagram
    add_frontend_architecture_slide(prs)

    # 7 — Backend APIs
    s7 = prs.slides.add_slide(prs.slide_layouts[6])
    add_table_slide(
        s7,
        "Backend — Route Handlers (src/app/api/)",
        ["Endpoint", "Purpose"],
        [
            ["POST /api/checkout", "Create Stripe Checkout Session with coupon + shipping"],
            ["POST /api/webhooks/stripe", "Verify signature; idempotent order on session.completed"],
            ["POST /api/orders/sync", "Fallback reconciliation when webhooks unavailable"],
            ["GET /api/orders · GET/POST …/[id]/cancel", "Authenticated order list, detail, cancel"],
            ["GET /api/products/search", "Rate-limited search (120/min/IP); Cache-Control 30s/SWR 60s"],
            ["POST /api/coupons/validate", "Server-side WELCOME10 / SAVE20 / FREESHIP rules"],
            ["GET/POST /api/auth/[...nextauth]", "NextAuth v5 session + OAuth callbacks"],
            ["POST /api/auth/signup · forgot/reset-password", "Registration and password recovery flows"],
        ],
    )

    # 8 — System architecture diagram
    add_system_architecture_slide(prs)

    # 9 — Checkout flow diagram
    add_checkout_flow_slide(prs)

    # 10 — Order flow diagram
    add_order_flow_slide(prs)

    # 11 — Search flow diagram
    add_search_flow_slide(prs)

    # 12 — Notifications & quality (visual layout)
    add_notifications_quality_slide(prs)

    # 13 — Deployment topology
    s13 = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_band(s13, "Deployment Topology")
    add_bullets(
        s13,
        [
            "Primary target: Vercel — vercel.json, iad1 region, security headers.",
            "Docker: node:22-alpine standalone Next.js output for portable runtimes.",
            "Kubernetes: Deployment (replicas: 2), Service, ConfigMap, Secret, probes on /.",
            "All third-party integrations are no-ops when env vars are unset — frictionless local dev.",
            "Recommended production path: Vercel + Postgres (Prisma/Drizzle) + Elastic Cloud trial/demo.",
            "Enterprise alternative: K8s + Docker with same codebase and observability agents.",
        ],
        size=16,
        space_after=Pt(18),
    )
    add_footer(s13)

    # 14 — Summary (visual)
    add_summary_slide(prs)

    # 15 — Thank you
    s14 = prs.slides.add_slide(prs.slide_layouts[6])
    bg14 = s14.background.fill
    bg14.solid()
    bg14.fore_color.rgb = DARK
    add_logo(s14, variant="light", top=Inches(0.55), left=Inches(0.7))
    accent14 = s14.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(3.2), SLIDE_W, Inches(0.08)
    )
    accent14.fill.solid()
    accent14.fill.fore_color.rgb = GREEN
    accent14.line.fill.background()
    thanks = s14.shapes.add_textbox(Inches(0.7), Inches(3.45), Inches(11.9), Inches(1.0))
    tr14 = thanks.text_frame.paragraphs[0]
    rr = tr14.add_run()
    rr.text = "Thank you"
    set_run(rr, size=44, bold=True, color=WHITE)
    add_bullets(
        s14,
        [
            "Pawan Gupta",
            "pawan.gupta@nagarro.com",
            "YCompany RIMSS — Tech Stack & Architecture",
        ],
        top=4.55,
        height=1.2,
        size=18,
        color=WHITE,
    )
    contact_box = s14.shapes[-1]
    for i, p in enumerate(contact_box.text_frame.paragraphs):
        for run in p.runs:
            set_run(run, size=18, color=GREEN if i == 0 else WHITE, bold=(i == 0))

    prs.save(OUT)
    print(f"Created {OUT}")


if __name__ == "__main__":
    build()
